"""Location-aware extraction + word-count sub-chunking.

Vendored from civ-kb's proven core (``aquarag.ingest.core``) as an extracted
library — Mode B does not import civ-kb at runtime (the vendoring decision, step
1). Each document is split at its natural structural unit (PDF page, PPTX slide,
DOCX heading, TXT/MD paragraph) before any word-count sub-chunking, so every
chunk keeps a human ``location`` ("page 2 [part 1/3]") for citation.

Two deliberate deviations from civ-kb (the mandatory corrections):

* **`.md` branch ADDED (carry-over #1).** civ-kb dispatches on pdf/docx/xlsx/
  pptx/txt only — it cannot read WDB's companion notes, where curated knowledge
  lives. ``_extract_md`` (≈ the ``.txt`` path) makes ``_context.md`` / ``_dict.md``
  / ``_about.md`` retrievable as passages. A conscious double-representation
  (note content is also in the graph) — but prose worth quoting.
* **`.csv` / `.xlsx` extractors DROPPED (carry-over #6).** Raw tables retrieve
  misleadingly (proof/FINDINGS.md Q2) and 100×-inflate the index; tabular data is
  answered by Mode C. There is no ``.csv`` branch here, by design — only prose.
"""

from __future__ import annotations

from pathlib import Path

from .model import CHUNK_OVERLAP, CHUNK_WORDS

# The prose formats Mode B indexes. Note the ABSENCE of .csv/.xlsx (carry-over #6)
# and the PRESENCE of .md (carry-over #1).
DOC_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


# ── PDF — per page via PyMuPDF ───────────────────────────────────────────────
def _extract_pdf(path: Path) -> list[tuple[str, str]]:
    """[(text, location), ...], one entry per page that has text."""
    try:
        import fitz  # PyMuPDF
    except Exception as e:  # pragma: no cover - import guard
        print(f"     [!] PyMuPDF unavailable: {e}")
        return []
    try:
        doc = fitz.open(str(path))
        out = []
        for i in range(len(doc)):
            text = doc[i].get_text("text").strip()
            if text:
                out.append((text, f"page {i + 1}"))
        doc.close()
        return out
    except Exception as e:  # pragma: no cover
        print(f"     [!] PyMuPDF error: {e}")
        return []


# ── PPTX — per slide ─────────────────────────────────────────────────────────
def _extract_pptx(path: Path) -> list[tuple[str, str]]:
    try:
        from pptx import Presentation
    except Exception as e:  # pragma: no cover
        print(f"     [!] python-pptx unavailable: {e}")
        return []
    try:
        prs = Presentation(str(path))
        out = []
        for n, slide in enumerate(prs.slides, start=1):
            parts = [
                para.text.strip()
                for shape in slide.shapes if shape.has_text_frame
                for para in shape.text_frame.paragraphs if para.text.strip()
            ]
            text = "\n".join(parts).strip()
            if text:
                out.append((text, f"slide {n}"))
        return out
    except Exception as e:  # pragma: no cover
        print(f"     [!] python-pptx error: {e}")
        return []


# ── DOCX — per heading/section ───────────────────────────────────────────────
def _extract_docx(path: Path) -> list[tuple[str, str]]:
    try:
        from docx import Document
    except Exception as e:  # pragma: no cover
        print(f"     [!] python-docx unavailable: {e}")
        return []
    try:
        doc = Document(str(path))
        sections, current = [], {"heading": None, "paras": []}
        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if style.startswith("Heading"):
                if current["paras"]:
                    sections.append(current)
                current = {"heading": text, "paras": []}
            else:
                current["paras"].append(text)
        for table in doc.tables:
            for row in table.rows:
                row_text = "  |  ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    current["paras"].append(row_text)
        if current["paras"]:
            sections.append(current)
        if not sections:
            return []
        has_headings = any(s["heading"] is not None for s in sections)
        out = []
        for i, s in enumerate(sections, start=1):
            text = "\n".join(s["paras"]).strip()
            if not text:
                continue
            loc = (s["heading"] or "Preamble") if has_headings else f"Block {i}"
            out.append((text, loc))
        return out
    except Exception as e:  # pragma: no cover
        print(f"     [!] python-docx error: {e}")
        return []


# ── TXT — per blank-line paragraph ───────────────────────────────────────────
def _extract_txt(path: Path) -> list[tuple[str, str]]:
    try:
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:  # pragma: no cover
        print(f"     [!] TXT read error: {e}")
        return []
    blocks = [b.strip() for b in text.split("\n\n") if b.strip()]
    return [(b, f"paragraph {i}") for i, b in enumerate(blocks, 1) if len(b.split()) > 5]


# ── MD — per blank-line block (carry-over #1) ────────────────────────────────
def _extract_md(path: Path) -> list[tuple[str, str]]:
    """[(text, location), ...] for a companion note (.md).

    ≈ the ``.txt`` path, but locations follow the note's Markdown ``##`` headings
    (e.g. "Summary", "Columns", "Related files") when present, so a cited note
    passage points at its section — closer to the §6 contract's ``note`` pointer.
    Blocks under no heading are labelled "block N".
    """
    try:
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:  # pragma: no cover
        print(f"     [!] MD read error: {e}")
        return []
    out: list[tuple[str, str]] = []
    section = "preamble"
    block_n = 0
    for raw in text.split("\n\n"):
        block = raw.strip()
        if not block:
            continue
        # update the current section label from a leading heading, if any
        first = block.splitlines()[0].strip()
        if first.startswith("#"):
            section = first.lstrip("#").strip() or section
        if len(block.split()) <= 3:  # a bare heading / divider — not a passage
            continue
        block_n += 1
        out.append((block, section))
    return out


# ── Dispatch ─────────────────────────────────────────────────────────────────
def extract_with_location(path) -> list[tuple[str, str]]:
    """[(text, location), ...] for any supported PROSE format (no .csv/.xlsx)."""
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(Path(path))
    if suffix == ".pptx":
        return _extract_pptx(Path(path))
    if suffix == ".docx":
        return _extract_docx(Path(path))
    if suffix == ".txt":
        return _extract_txt(Path(path))
    if suffix == ".md":
        return _extract_md(Path(path))
    return []


# ── Sub-chunker — preserves the location label ───────────────────────────────
def sub_chunk(text: str, location: str,
              chunk_words: int = CHUNK_WORDS, overlap: int = CHUNK_OVERLAP) -> list[tuple[str, str]]:
    """Split an oversized natural unit into word-count chunks (vendored verbatim).

    Location becomes e.g. "page 14 [part 1/3]" so provenance is never lost.
    """
    words = text.split()
    if len(words) <= chunk_words:
        return [(text, location)]
    raw, i = [], 0
    while i < len(words):
        chunk = " ".join(words[i:i + chunk_words])
        if len(chunk.strip()) > 80:
            raw.append(chunk)
        i += chunk_words - overlap
    total = len(raw)
    if total <= 1:
        return [(raw[0] if raw else text, location)]
    return [(c, f"{location} [part {idx + 1}/{total}]") for idx, c in enumerate(raw)]
